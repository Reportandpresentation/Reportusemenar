import streamlit as st
import google.generativeai as genai
from docx import Document
from pptx import Presentation
from io import BytesIO

# --- 1. ڕێکخستنا API و پاسوۆردێ ئەدمینی ---
try:
    API_KEY = st.secrets["AIzaSyCImqsHi-DyswsFvbY2q32qROLZVaE1j-o"]
except:
    API_KEY = "لێرە_کلیلێ_خۆ_دانە"

# لێرە پاسوۆردێ خۆ دیار بکە کو دێ دەیە قوتابیان
ADMIN_PASSWORD = "REPORT_2024" 

genai.configure(api_key=API_KEY)

# --- 2. دیزاینا CSS ---
st.markdown("""
    <style>
    .stApp { background-color: #f4f7f6; }
    .price-tag {
        background: #002b5b; color: #d4af37; padding: 15px;
        border-radius: 10px; text-align: center; font-weight: bold; font-size: 22px;
    }
    .payment-box {
        background: #ffffff; padding: 25px; border-radius: 15px;
        border: 2px solid #d4af37; margin-top: 20px;
    }
    </style>
    """, unsafe_allow_html=True)

st.markdown('<div class="price-tag">بهایێ هەر ڕاپۆرت یان سمینارەکێ تنێ 5,000 دینارە (Zain Card)</div>', unsafe_allow_html=True)
st.title("ناڤەندا زیرەک بۆ خزمەتێن زانستی 🎓")

# --- 3. فۆڕمێ داخوازیا ڕاپۆرتێ ---
with st.expander("📝 لێرە دەستپێ بکە و زانیاریان تژی بکە", expanded=True):
    title = st.text_input("ناڤنیشانێ بابەت (Title):")
    category = st.selectbox("جۆرێ کارێ:", ["ڕاپۆرت", "سمینار (PowerPoint)"])
    desc = st.text_area("وەسفەکا کورت ل سەر بابەتێ:")
    lang = st.selectbox("زمان:", ["Kurdish", "Arabic", "English"])
    generate_btn = st.button("دروست بکە و پێشاندە ✨")

# --- 4. پرۆسێسا دروستکرنێ ---
if generate_btn:
    if not title or not desc:
        st.error("تکایە هەمی خانەیان تژی بکە!")
    else:
        with st.spinner('ل حالەتێ دروستکرنا نموونەیێ دایە...'):
            model = genai.GenerativeModel('gemini-1.5-flash')
            prompt = f"Create a high-quality academic {category} about {title} in {lang}. Description: {desc}. Include references."
            response = model.generate_content(prompt)
            st.session_state['full_content'] = response.text
            st.session_state['work_title'] = title
            st.session_state['work_type'] = category
            st.session_state['step'] = 'payment'

# --- 5. قوناغا پارەدان و داونلۆدکرنێ ---
if 'step' in st.session_state:
    st.divider()
    st.subheader("📊 پێشاندانا کورت (Preview)")
    st.info(st.session_state['full_content'][:400] + "...")
    st.warning("بۆ داونلۆدکرنا فایلا تەمام، پێدڤیە تو پاسوۆردێ چالاککرنێ بنڤێسی.")

    # باکسێ پارەدانێ
    with st.container():
        st.markdown('<div class="payment-box">', unsafe_allow_html=True)
        st.markdown("### 💳 ڕێکارێن وەرگرتنا فایلی:")
        st.write("1. کۆدێ کارتا زەین (5,000) ب وێنە ڤە بفرێکه بۆ واتساپی.")
        st.write("2. پشتی پشتڕاستکرنێ، ئەم دێ پاسوۆردێ داونلۆدێ بۆ تە فرێکەین.")
        
        # دوگمەیا واتساپی
        wa_msg = f"سڵاو، من داخوازیەکا {st.session_state['work_type']} کری ل سەر بابەتێ ({st.session_state['work_title']}). تکایە پاسوۆردێ داونلۆدێ بنێرە."
        wa_link = f"https://wa.me/9647508015653?text={wa_msg.replace(' ', '%20')}"
        st.markdown(f'<a href="{wa_link}" target="_blank" style="background-color: #25D366; color: white; padding: 12px; text-decoration: none; border-radius: 8px; display: inline-block;">Click to send Card via WhatsApp 💬</a>', unsafe_allow_html=True)
        
        st.markdown("---")
        
        # پشکا پاسوۆردی
        user_pwd = st.text_input("🔑 پاسوۆردێ وەرگرتی لێرە بنڤێسه:", type="password")
        
        if user_pwd == ADMIN_PASSWORD:
            st.success("پاسوۆرد یێ درستە! نوکە دشێی فایلی داونلۆد بکەی.")
            
            # دروستکرنا فایلا داونلۆدێ
            if st.session_state['work_type'] == "ڕاپۆرت":
                doc = Document()
                doc.add_heading(st.session_state['work_title'], 0)
                doc.add_paragraph(st.session_state['full_content'])
                buf = BytesIO()
                doc.save(buf)
                st.download_button("📥 داونلۆدکرنا فایلا Word", data=buf.getvalue(), file_name=f"{st.session_state['work_title']}.docx")
            else:
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = st.session_state['work_title']
                buf = BytesIO()
                prs.save(buf)
                st.download_button("📥 داونلۆدکرنا فایلا PowerPoint", data=buf.getvalue(), file_name=f"{st.session_state['work_title']}.pptx")
        elif user_pwd != "":
            st.error("پاسوۆرد یێ خەلەتە! تکایە پەیوەندیێ ب ئەدمینی بکە.")
        st.markdown('</div>', unsafe_allow_html=True)
